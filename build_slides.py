#!/usr/bin/env python3
"""
build_slides.py — render the SF Civic Art Map deck as a PPTX with the same
typographic register as the HTML version. Upload to Google Drive and right-
click → Open with Google Slides to convert.

Fonts referenced (all available in Google Slides via Google Fonts):
  - Cormorant Garamond (serif headlines)
  - Inter (sans body)
  - JetBrains Mono (mono eyebrows + labels)

Colors match the live deck:
  bg=#FBF8F2  bg-alt=#EDF1F8  bg-panel=#E4ECF7
  text=#1C2638  dim=#5C6577  headline=#0E1B36  accent=#1F3D88
"""

from __future__ import annotations
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
from lxml import etree

# ─── Palette ───────────────────────────────────────────────────────────────
BG          = RGBColor(0xFB, 0xF8, 0xF2)  # cream
BG_ALT      = RGBColor(0xED, 0xF1, 0xF8)
BG_PANEL    = RGBColor(0xE4, 0xEC, 0xF7)
TEXT        = RGBColor(0x1C, 0x26, 0x38)
TEXT_DIM    = RGBColor(0x5C, 0x65, 0x77)
HEADLINE    = RGBColor(0x0E, 0x1B, 0x36)
ACCENT      = RGBColor(0x1F, 0x3D, 0x88)
ACCENT_SOFT = RGBColor(0x4E, 0x72, 0xB8)
RULE        = RGBColor(0xD6, 0xDE, 0xEA)

# ─── Fonts ─────────────────────────────────────────────────────────────────
SERIF = "Cormorant Garamond"
SANS  = "Inter"
MONO  = "JetBrains Mono"

# ─── Geometry (16:9) ──────────────────────────────────────────────────────
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)
PAD_X   = Inches(0.85)
PAD_TOP = Inches(0.85)
PAD_BOT = Inches(0.55)
NAV_H   = Inches(0.50)  # brand strip at top of every slide
CONTENT_W = SLIDE_W - 2 * PAD_X

# ─── Type scale (mirrors the web's clamp() values, in pt) ────────────────
H1_SIZE       = 48
H2_SIZE       = 32
H3_SIZE       = 20
DEK_SIZE      = 15
LEAD_SIZE     = 14
BODY_SIZE     = 13
EYEBROW_SIZE  = 11
BYLINE_SIZE   = 10
THESIS_SIZE   = 22
CARD_NAME_SIZE = 19
PHASE_NUM_SIZE = 46
CLOSING_SIZE  = 96


# ───────────────────────────────────────────────────────────────────────────
# Helpers
# ───────────────────────────────────────────────────────────────────────────

def set_bg(slide, color: RGBColor):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_text_box(slide, x, y, w, h, *, anchor=MSO_ANCHOR.TOP):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    # Clear default empty paragraph spacing
    return tf


def run(p, text, *, font=SANS, size=12, bold=False, italic=False, color=TEXT,
        space_after=Pt(0), align=PP_ALIGN.LEFT, char_spacing=None):
    if p.text:  # If paragraph already has text, add a new run
        r = p.add_run()
    else:
        r = p.add_run()
    r.text = text
    r.font.name = font
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.italic = italic
    r.font.color.rgb = color
    if char_spacing is not None:
        # Apply character spacing via XML (python-pptx doesn't expose this)
        rPr = r._r.get_or_add_rPr()
        rPr.set("spc", str(int(char_spacing * 100)))
    p.alignment = align
    if space_after:
        p.space_after = space_after
    return r


def add_paragraph(tf, *, first=False):
    if first and not tf.paragraphs[0].text:
        return tf.paragraphs[0]
    return tf.add_paragraph()


def eyebrow_para(tf, text, *, first=False):
    p = add_paragraph(tf, first=first)
    run(p, text.upper(), font=MONO, size=11, bold=True, color=ACCENT,
        char_spacing=2.4, space_after=Pt(14))
    return p


def h1_para(tf, text, *, first=False, size=34):
    p = add_paragraph(tf, first=first)
    run(p, text, font=SERIF, size=size, color=HEADLINE, space_after=Pt(18))
    return p


def h2_para(tf, text, *, first=False, size=26):
    p = add_paragraph(tf, first=first)
    run(p, text, font=SERIF, size=size, color=HEADLINE, space_after=Pt(14))
    return p


def body_para(tf, text, *, first=False, size=13, color=TEXT, space=10, align=PP_ALIGN.LEFT):
    p = add_paragraph(tf, first=first)
    run(p, text, font=SANS, size=size, color=color, space_after=Pt(space), align=align)
    return p


def label_para(tf, text, *, first=False, size=9, color=ACCENT, space=4):
    p = add_paragraph(tf, first=first)
    run(p, text.upper(), font=MONO, size=size, bold=True, color=color,
        char_spacing=2.0, space_after=Pt(space))
    return p


def add_rect(slide, x, y, w, h, *, fill=None, line=None):
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    if fill is None:
        shp.fill.background()
    else:
        shp.fill.solid()
        shp.fill.fore_color.rgb = fill
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line
        shp.line.width = Pt(0.75)
    shp.shadow.inherit = False
    return shp


def add_thin_rule(slide, x, y, w, color=RULE, height=Pt(0.75)):
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, Emu(int(height)))
    shp.fill.solid()
    shp.fill.fore_color.rgb = color
    shp.line.fill.background()
    shp.shadow.inherit = False
    return shp


def add_brand_strip(slide):
    """Top-of-slide brand mark mirroring the website's nav."""
    tf = add_text_box(slide, PAD_X, Inches(0.30), Inches(8), Inches(0.35),
                      anchor=MSO_ANCHOR.MIDDLE)
    p = tf.paragraphs[0]
    run(p, "SAN FRANCISCO · CIVIC ART MAP", font=MONO, size=11, bold=True,
        color=HEADLINE, char_spacing=3.0)
    add_thin_rule(slide, 0, Inches(0.82), SLIDE_W, color=RULE)


# Content area starts below the brand strip
CONTENT_TOP = Inches(1.2)


# ───────────────────────────────────────────────────────────────────────────
# Slides
# ───────────────────────────────────────────────────────────────────────────

def slide_hero(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    set_bg(s, BG)
    add_brand_strip(s)

    # Eyebrow
    eb_tf = add_text_box(s, PAD_X, Inches(1.45), CONTENT_W, Inches(0.32))
    p = eb_tf.paragraphs[0]
    run(p, "A PROPOSAL FOR ANDY AND DEBORAH RAPPAPORT", font=MONO,
        size=EYEBROW_SIZE, bold=True, color=ACCENT, char_spacing=3.5)

    # H1 — large serif with italic-accent on "civic art utility"
    h1_tf = add_text_box(s, PAD_X, Inches(1.95), CONTENT_W, Inches(3.2))
    p = h1_tf.paragraphs[0]
    p.line_spacing = 1.06
    run(p, "A ", font=SERIF, size=H1_SIZE, color=HEADLINE)
    run(p, "civic art utility", font=SERIF, size=H1_SIZE, italic=True, color=ACCENT)
    run(p, " for San Francisco—enriching, beautifully made, "
          "and built as a laboratory for the next generation of cultural "
          "infrastructure.",
        font=SERIF, size=H1_SIZE, color=HEADLINE)

    # Dek
    dek_tf = add_text_box(s, PAD_X, Inches(5.30), CONTENT_W, Inches(1.30))
    p = dek_tf.paragraphs[0]
    p.line_spacing = 1.42
    run(p,
        "San Francisco has one of the greatest art ecosystems in the country "
        "and one of the hardest to navigate. The opportunity is a single, "
        "comprehensive cultural utility for the city—built with care, run in "
        "the public interest, and designed as a shared working environment "
        "for the people, institutions, and technologies that could make San "
        "Francisco the most rewarding place in the world to encounter art.",
        font=SANS, size=DEK_SIZE, color=TEXT)

    # Byline pinned near bottom
    by_tf = add_text_box(s, PAD_X, Inches(6.80), CONTENT_W, Inches(0.35),
                         anchor=MSO_ANCHOR.MIDDLE)
    p = by_tf.paragraphs[0]
    run(p, "PREPARED BY ", font=MONO, size=BYLINE_SIZE, color=TEXT_DIM, char_spacing=2.5)
    run(p, "ANDREW GOLDSTEIN", font=MONO, size=BYLINE_SIZE, bold=True, color=ACCENT,
        char_spacing=2.5)
    run(p, " · DEDALUS MEDIA ARCHITECTS · MAY 2026", font=MONO, size=BYLINE_SIZE,
        color=TEXT_DIM, char_spacing=2.5)


def slide_vision(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s, BG)
    add_brand_strip(s)

    # Eyebrow + H2
    head_tf = add_text_box(s, PAD_X, CONTENT_TOP, CONTENT_W, Inches(2.1))
    eyebrow_para(head_tf, "The Vision", first=True)
    p = head_tf.add_paragraph()
    p.line_spacing = 1.14
    run(p,
        "An art map for San Francisco that feels like the city itself—",
        font=SERIF, size=H2_SIZE, color=HEADLINE)
    run(p, "inspiring, generous, delightful, and unafraid of beauty",
        font=SERIF, size=H2_SIZE, italic=True, color=HEADLINE)
    run(p, ".", font=SERIF, size=H2_SIZE, color=HEADLINE)

    # Body paragraphs
    body_tf = add_text_box(s, PAD_X, Inches(3.6), CONTENT_W, Inches(2.4))
    p = body_tf.paragraphs[0]
    p.line_spacing = 1.5
    run(p,
        "For all the energy in the city's cultural life, there is no "
        "comprehensive way to find and engage with it. Museums, galleries, "
        "pop-ups, artist studios, public art, AR and digital-art activations "
        "all coexist, but none are legible to one another—let alone to the "
        "local audience or curious visitor.",
        font=SANS, size=BODY_SIZE, color=TEXT)
    p = body_tf.add_paragraph()
    p.line_spacing = 1.5
    p.space_before = Pt(10)
    run(p,
        "This is a market-design problem. The traditional venues—galleries, "
        "museums, nonprofits—are under structural pressure. Vital newcomers"
        "—digital artists, immersive venues, AI-native creators—are operating "
        "without an infrastructure. Critics who previously served as wayfinding "
        "mechanisms have either vanished or been fragmented across "
        "algorithmic channels.",
        font=SANS, size=BODY_SIZE, color=TEXT)

    # Thesis — tall thin accent bar (analog of the web's arrow indicator)
    # + larger serif text to its right.
    thesis_y = Inches(6.10)
    thesis_h = Inches(0.95)
    bar_w = Inches(0.09)
    add_rect(s, PAD_X, thesis_y, bar_w, thesis_h, fill=ACCENT)
    thesis_tf = add_text_box(s, PAD_X + bar_w + Inches(0.30), thesis_y,
                             CONTENT_W - bar_w - Inches(0.30), thesis_h,
                             anchor=MSO_ANCHOR.MIDDLE)
    p = thesis_tf.paragraphs[0]
    p.line_spacing = 1.30
    run(p,
        "What's needed is not another platform that competes with these "
        "ecosystem participants, but a piece of civic infrastructure that "
        "connects them to one another and surfaces them to the audiences "
        "they need to thrive.",
        font=SERIF, size=THESIS_SIZE, color=HEADLINE)


def slide_commitments(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s, BG_ALT)
    add_brand_strip(s)

    head_tf = add_text_box(s, PAD_X, CONTENT_TOP, CONTENT_W, Inches(1.8))
    eyebrow_para(head_tf, "Three commitments shape it", first=True)
    h2_para(head_tf,
            "A new category of cultural utility, designed by what it refuses to be.",
            size=H2_SIZE)

    rows = [
        ("Civic, not commercial.",
         "Operated in the public interest, with editorial independence and a sustaining model that doesn't depend on extracting value from artists or galleries. It complements—never competes with—the institutions whose work it surfaces."),
        ("Beautifully made, not merely functional.",
         "The interface holds itself to a standard usually reserved for cultural institutions: the care of a national-park map, the warmth of a hospitable host, the clarity of a great editorial product. San Francisco deserves nothing less than the most thoughtfully designed cultural utility any American city has built."),
        ("A laboratory, not a finished product.",
         "A working environment where the next generation of interpretive technologies—AI-led understanding, spatial audio, AR overlays, autonomous-vehicle cultural routing—can be deployed, tested, and refined in real public use, with the active participation of the companies that are building them."),
    ]
    y = Inches(3.4)
    row_h = Inches(1.25)
    term_w = Inches(3.6)
    gap = Inches(0.5)
    for i, (term, desc) in enumerate(rows):
        ty = y + row_h * i
        tf_term = add_text_box(s, PAD_X, ty + Inches(0.05), term_w, row_h, anchor=MSO_ANCHOR.TOP)
        p = tf_term.paragraphs[0]
        p.line_spacing = 1.18
        run(p, term, font=SERIF, size=22, italic=True, color=ACCENT)

        tf_desc = add_text_box(s, PAD_X + term_w + gap, ty + Inches(0.05),
                               CONTENT_W - term_w - gap, row_h)
        p = tf_desc.paragraphs[0]
        p.line_spacing = 1.45
        run(p, desc, font=SANS, size=BODY_SIZE, color=TEXT)

        if i < len(rows) - 1:
            add_thin_rule(s, PAD_X, ty + row_h - Inches(0.05), CONTENT_W)


def slide_why_now(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s, BG)
    add_brand_strip(s)

    head_tf = add_text_box(s, PAD_X, CONTENT_TOP, CONTENT_W, Inches(1.8))
    eyebrow_para(head_tf, "Why Now", first=True)
    h2_para(head_tf,
            "Three currents are converging on this exact problem at this exact moment.",
            size=H2_SIZE)

    columns = [
        ("The cultural moment", [
            "San Francisco's gallery and institutional ecosystem has been reshaped by a decade of rent pressure, the long tail of the pandemic, and the recent transitions at SFAI, ICA SF, and elsewhere.",
            "The institutions that remain are working harder than ever to reach audiences while operating in relative isolation from one another.",
            "The city's cultural map is fragmenting at exactly the moment when Mayor Lurie has made cultural revitalization a stated priority.",
        ]),
        ("The technology moment", [
            "For the first time, the technology to actually unlock the art experience exists. AI can read a wall label, a press release, an artist's Q&A, and a curator's catalogue together—and offer the visitor a personalized way in.",
            "AR glasses and spatial-audio earbuds are about to make on-site interpretation effortless. Autonomous vehicles will reshape how visitors move through the city.",
            "The next twenty-four months are when these capabilities go from emerging to ambient, and the city that builds the showcase first will set the template.",
        ]),
        ("The civic moment", [
            "San Francisco's tech leadership is under pressure to make the city a place worth living in—not just as philanthropy, but because the quality of civic life is foundational to attracting and keeping top talent.",
            "Voices like Patrick Collison's call for new aesthetics and Packy McCormick's appeal for modern magnificenza signal a hunger across the tech community for positive, generous expressions of the technological age.",
            "The right civic initiative, designed with discipline, can convene that goodwill into something durable, useful, and a credit to everyone involved.",
        ]),
    ]
    col_w = (CONTENT_W - Inches(0.8)) / 3
    y = Inches(3.45)
    for i, (head, bullets) in enumerate(columns):
        x = PAD_X + (col_w + Inches(0.4)) * i
        # accent rule at top
        add_rect(s, x, y, col_w, Pt(2.5), fill=ACCENT)
        ctf = add_text_box(s, x, y + Inches(0.14), col_w, Inches(3.5))
        p = ctf.paragraphs[0]
        p.line_spacing = 1.15
        run(p, head, font=SERIF, size=24, color=HEADLINE, space_after=Pt(11))
        for b in bullets:
            bp = ctf.add_paragraph()
            bp.line_spacing = 1.42
            bp.space_after = Pt(8)
            run(bp, "•  ", font=SANS, size=12, bold=True, color=ACCENT)
            run(bp, b, font=SANS, size=12, color=TEXT)


def slide_who(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s, BG_ALT)
    add_brand_strip(s)

    head_tf = add_text_box(s, PAD_X, CONTENT_TOP, CONTENT_W, Inches(1.8))
    eyebrow_para(head_tf, "Who it's for", first=True)
    h2_para(head_tf,
            "The map serves the full population of people who want to encounter "
            "art in San Francisco—as well as those who don't know they want to yet.",
            size=H2_SIZE)

    cards = [
        ("Saturday afternoon", "The local out exploring",
         "Lives in the city, knows enough to be choosy, often runs out of ideas when trying to figure out what to do."),
        ("On the hunt", "The collector hitting the circuit",
         "Lives for openings, studio visits, and finding the next world-changing artist."),
        ("Weekend family", "The family with kids",
         "Wants culture as part of a weekend, not as an obligation. Friendly, age-aware, reachable."),
        ("Ambitious learner", "The student on a budget",
         "Curious, idea-hungry, stretching every cultural dollar. What's free, near transit, eye-sharpening."),
        ("Training taste", "The tech worker seeking inspo",
         "Building the next generation of tools and looking for cultural input that sharpens taste."),
        ("Visitor with a weekend", "The tourist on the town",
         "A few hours, doesn't know the neighborhoods, wants the best of SF on a walkable itinerary."),
    ]
    y0 = Inches(3.55)
    cols = 3
    gap = Inches(0.22)
    card_w = (CONTENT_W - gap * (cols - 1)) / cols
    card_h = Inches(1.75)
    for i, (role, name, desc) in enumerate(cards):
        cx = PAD_X + (card_w + gap) * (i % cols)
        cy = y0 + (card_h + gap) * (i // cols)
        add_rect(s, cx, cy, card_w, card_h, fill=BG, line=RULE)
        ctf = add_text_box(s, cx + Inches(0.20), cy + Inches(0.16),
                           card_w - Inches(0.40), card_h - Inches(0.32))
        label_para(ctf, role, first=True, space=4)
        p = ctf.add_paragraph()
        p.line_spacing = 1.14
        run(p, name, font=SERIF, size=CARD_NAME_SIZE, color=HEADLINE,
            space_after=Pt(6))
        p = ctf.add_paragraph()
        p.line_spacing = 1.40
        run(p, desc, font=SANS, size=11, color=TEXT)


def slide_experience(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s, BG)
    add_brand_strip(s)

    head_tf = add_text_box(s, PAD_X, CONTENT_TOP, CONTENT_W, Inches(2.3))
    eyebrow_para(head_tf, "The Experience", first=True)
    h2_para(head_tf,
            "A single user journey shapes the whole product: prepare, encounter, reward.",
            size=H2_SIZE)
    p = head_tf.add_paragraph()
    p.line_spacing = 1.42
    p.space_before = Pt(4)
    run(p,
        "The map's job is to use digital tools to funnel people into the "
        "real—and then use digital tools to augment their experience. The "
        "interpretive layer is built on a three-part choreography that "
        "lowers the barrier to understanding and elevates the visitor's intent.",
        font=SANS, size=LEAD_SIZE, color=TEXT)

    steps = [
        ("I.", "Prepare",
         "Before the visit, the map gives the visitor the priming context that makes the encounter resonant—curator-vetted, AI-assisted, voiced in the institution's register."),
        ("II.", "Encounter",
         "During the visit, the map is a helpful expert companion—audio guides tuned to one's preference, a wayfinding layer that respects the institution's authority."),
        ("III.", "Reward",
         "Virtuous gamification—streaks, badges—turns completed itineraries and repeat visits into status. Light-lift reviews become attention-directing signals."),
    ]
    y = Inches(4.05)
    col_w = (CONTENT_W - Inches(0.8)) / 3
    for i, (num, name, desc) in enumerate(steps):
        x = PAD_X + (col_w + Inches(0.4)) * i
        ctf = add_text_box(s, x, y, col_w, Inches(1.8))
        p = ctf.paragraphs[0]
        run(p, num, font=MONO, size=13, bold=True, color=ACCENT, space_after=Pt(4))
        p = ctf.add_paragraph()
        p.line_spacing = 1.14
        run(p, name, font=SERIF, size=26, color=HEADLINE, space_after=Pt(6))
        p = ctf.add_paragraph()
        p.line_spacing = 1.42
        run(p, desc, font=SANS, size=11.5, color=TEXT)

    # Lenses chips — sized to their text content, like the web pill widths
    chips = ["Time", "Distance", "Mood", "Preferences", "Network"]
    chip_h = Inches(0.30)
    gap_x = Inches(0.10)
    # JetBrains Mono at 10pt: ~5.7pt per char; +14pt horizontal padding
    chip_widths = [Inches((len(c) * 5.7 + 16) / 72) for c in chips]
    total_w = sum(chip_widths, Inches(0)) + gap_x * (len(chips) - 1)
    chip_y = Inches(6.10)
    cx = (SLIDE_W - total_w) / 2
    for chip, cw in zip(chips, chip_widths):
        shp = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, cx, chip_y, cw, chip_h)
        shp.adjustments[0] = 0.5
        shp.fill.solid()
        shp.fill.fore_color.rgb = BG
        shp.line.color.rgb = RULE
        shp.line.width = Pt(0.5)
        shp.shadow.inherit = False
        ctf = shp.text_frame
        ctf.margin_left = Pt(4); ctf.margin_right = Pt(4)
        ctf.margin_top = Pt(0); ctf.margin_bottom = Pt(0)
        ctf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = ctf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run(p, chip, font=MONO, size=10, color=TEXT)
        cx += cw + gap_x

    # Caption under the lenses
    cap_y = chip_y + chip_h + Inches(0.12)
    cap_tf = add_text_box(s, PAD_X, cap_y, CONTENT_W, Inches(0.3))
    p = cap_tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run(p, "Composable interface primitives that combine into emergent journeys.",
        font=SANS, size=11, italic=True, color=TEXT_DIM)


def slide_unlock(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s, BG_ALT)
    add_brand_strip(s)

    head_tf = add_text_box(s, PAD_X, CONTENT_TOP, CONTENT_W, Inches(2.5))
    eyebrow_para(head_tf, "The Unlock", first=True)
    h2_para(head_tf,
            "San Francisco's tech leaders are art people. The map gives them "
            "a way to contribute that isn't writing a check.",
            size=H2_SIZE)
    p = head_tf.add_paragraph()
    p.line_spacing = 1.42
    p.space_before = Pt(4)
    run(p,
        "The founders building the next generation of technology are often "
        "deeply interested in art. They collect. They support institutions. "
        "They sit on cultural boards. They are surrounded by colleagues who "
        "would happily contribute talent, infrastructure, and access to a "
        "worthy civic project—if the project were ambitious enough to deserve it.",
        font=SANS, size=BODY_SIZE, color=TEXT)

    modes = [
        ("As civic collaborators",
         "The map becomes a venue where their interest in art finds civic expression. Their teams can volunteer design and engineering. Their companies can lend infrastructure. The act of asking for capability rather than capital changes the whole register of the conversation."),
        ("As laboratory partners",
         "The map is a working environment for the technologies these companies are actively developing. AR glasses, spatial-audio interfaces, autonomous-vehicle services, AI-led interpretation—all need a curated, well-run real-world deployment surface. The map is that surface."),
        ("As founding stakeholders",
         "The map is designed to scale, but not extractively—a civic utility with a sustaining business model. The companies who help shape it here become the founding stakeholders of a new category of cultural infrastructure."),
    ]
    y = Inches(4.65)
    col_w = (CONTENT_W - Inches(0.5)) / 3
    panel_h = Inches(1.55)
    for i, (head, desc) in enumerate(modes):
        x = PAD_X + (col_w + Inches(0.25)) * i
        add_rect(s, x, y, col_w, panel_h, fill=BG_PANEL)
        add_rect(s, x, y, col_w, Pt(3), fill=ACCENT)
        ctf = add_text_box(s, x + Inches(0.18), y + Inches(0.14),
                           col_w - Inches(0.36), panel_h - Inches(0.28))
        p = ctf.paragraphs[0]
        p.line_spacing = 1.18
        run(p, head, font=SERIF, size=20, color=HEADLINE, space_after=Pt(6))
        p = ctf.add_paragraph()
        p.line_spacing = 1.40
        run(p, desc, font=SANS, size=11, color=TEXT)

    # Banner
    by = Inches(6.40)
    bh = Inches(0.65)
    bshp = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, PAD_X, by, CONTENT_W, bh)
    bshp.fill.background()
    bshp.line.color.rgb = ACCENT
    bshp.line.width = Pt(0.75)
    bshp.shadow.inherit = False
    btf = bshp.text_frame
    btf.margin_left = Pt(14); btf.margin_right = Pt(14)
    btf.margin_top = Pt(4); btf.margin_bottom = Pt(4)
    btf.vertical_anchor = MSO_ANCHOR.MIDDLE
    btf.word_wrap = True
    p = btf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    p.line_spacing = 1.25
    run(p,
        "Tech companies don't want to be philanthropists; they want to be "
        "participants. The map gives them a participant role in their city's "
        "cultural life—and a laboratory for the cultural products they're "
        "already building.",
        font=SERIF, size=16, color=HEADLINE)


def slide_how(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s, BG)
    add_brand_strip(s)

    head_tf = add_text_box(s, PAD_X, CONTENT_TOP, CONTENT_W, Inches(1.8))
    eyebrow_para(head_tf, "How we build it", first=True)
    h2_para(head_tf,
            "A three-phase arc, with each phase committing only what the prior phase has earned.",
            size=H2_SIZE)

    phases = [
        ("01", "Discovery · 4–6 weeks", "Listen, scope, and lay the foundation.",
         "A focused engagement that interviews the stakeholders, surveys the peer landscape, identifies the partner-integration vectors, and produces a draft scope-and-plan document for Phase 2. Light-footprint, decisive, designed so anyone reading it can see exactly what they would be funding and why."),
        ("02", "Prototype · 8–12 weeks", "Vibe-code a working version. Put it in real hands.",
         "A working web application with a curated subset of the city's institutions, the first version of the interpretation layer, and the first integrations with the most willing tech partners. Tested with real users, iterated weekly. The deliverable is not a deck; it is a thing people use."),
        ("03", "Launch + Steward · Open structure", "Ship the public utility. Make it last.",
         "The full civic launch—operating structure, sustaining model, governance, and the team to run it. The shape of Phase 3 is what Phases 1 and 2 are designed to figure out."),
    ]
    y0 = Inches(3.55)
    row_h = Inches(1.20)
    num_w = Inches(1.05)
    gap = Inches(0.35)
    for i, (num, label, head, desc) in enumerate(phases):
        ry = y0 + row_h * i
        num_tf = add_text_box(s, PAD_X, ry, num_w, row_h, anchor=MSO_ANCHOR.TOP)
        p = num_tf.paragraphs[0]
        run(p, num, font=SERIF, size=PHASE_NUM_SIZE, color=ACCENT)

        body_x = PAD_X + num_w + gap
        body_w = CONTENT_W - num_w - gap
        b_tf = add_text_box(s, body_x, ry + Inches(0.05), body_w, row_h)
        p = b_tf.paragraphs[0]
        run(p, label.upper(), font=MONO, size=10, bold=True, color=TEXT_DIM,
            char_spacing=2.4, space_after=Pt(3))
        p = b_tf.add_paragraph()
        p.line_spacing = 1.18
        run(p, head, font=SERIF, size=20, color=HEADLINE, space_after=Pt(4))
        p = b_tf.add_paragraph()
        p.line_spacing = 1.42
        run(p, desc, font=SANS, size=11.5, color=TEXT)

        if i < len(phases) - 1:
            add_thin_rule(s, PAD_X, ry + row_h - Inches(0.05), CONTENT_W)


def slide_phase1(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s, BG_ALT)
    add_brand_strip(s)

    head_tf = add_text_box(s, PAD_X, CONTENT_TOP, CONTENT_W, Inches(2.5))
    eyebrow_para(head_tf, "Phase 1 · The ask", first=True)
    h2_para(head_tf,
            "A focused discovery sprint—enough to know exactly what we're "
            "building and what it will take.",
            size=H2_SIZE)
    p = head_tf.add_paragraph()
    p.line_spacing = 1.42
    p.space_before = Pt(4)
    run(p,
        "Phase 1 is option-pricing. The proposal is to fund a structured "
        "4–6 week discovery sprint that convenes the right stakeholders, "
        "surfaces what the city actually needs, and produces a draft scope "
        "and Phase 2 proposal you can react to with full information.",
        font=SANS, size=BODY_SIZE, color=TEXT)

    # Content-sized row heights — Scope needs more room for its 6 bullets,
    # the single-line rows hug their content. Last-row gets a tiny bit extra
    # so its 2-line description doesn't crowd the frame bottom.
    rows = [
        ("Scope", [
            "20–30 stakeholder interviews across institutions, galleries, artists, civic, tech",
            "Peer landscape and best-practice synthesis",
            "Product scope for the Phase 2 prototype",
            "Partnership architecture and sequence",
            "Governance and sustaining-model options",
            "Phase 2 proposal with timeline, team, budget",
        ], Inches(1.10)),
        ("Timeline", ["Four to six weeks from kickoff to delivered document."], Inches(0.30)),
        ("Fee", ["$75,000 fixed, plus travel and expenses."], Inches(0.30)),
        ("Payment", ["Half upon commencement, half upon completion."], Inches(0.30)),
        ("What you provide", ["Vision alignment, input, and ideas; funding; introductions to art and tech stakeholders."], Inches(0.40)),
        ("What happens at the end", ["You have a document you can read, react to, and decide on without commitment. If Phase 2 is the right move, we proceed."], Inches(0.50)),
    ]

    # Compute total frame height from row heights + rules + padding
    inner_pad = Inches(0.18)
    rule_thick = Pt(0.5)
    rows_total = sum(h for _, _, h in rows)
    rules_total = rule_thick * (len(rows) - 1)
    detail_h = rows_total + rules_total + 2 * inner_pad
    detail_y = Inches(3.85)

    # Frame
    add_rect(s, PAD_X, detail_y, CONTENT_W, detail_h, fill=BG, line=RULE)

    inner_x = PAD_X + inner_pad
    inner_w = CONTENT_W - 2 * inner_pad
    label_w = Inches(2.0)
    body_x = inner_x + label_w + Inches(0.25)
    body_w = inner_w - label_w - Inches(0.25)

    # Lay out rows top-down with content-driven heights
    cur_y = detail_y + inner_pad
    for i, (label, items, h) in enumerate(rows):
        # Label
        l_tf = add_text_box(s, inner_x, cur_y, label_w, h, anchor=MSO_ANCHOR.TOP)
        p = l_tf.paragraphs[0]
        run(p, label.upper(), font=MONO, size=9, bold=True, color=ACCENT,
            char_spacing=2.0)

        # Body
        b_tf = add_text_box(s, body_x, cur_y, body_w, h, anchor=MSO_ANCHOR.TOP)
        if len(items) == 1:
            p = b_tf.paragraphs[0]
            run(p, items[0], font=SANS, size=10.5, color=TEXT)
        else:
            for j, it in enumerate(items):
                p = (b_tf.paragraphs[0] if j == 0 else b_tf.add_paragraph())
                run(p, "•  ", font=SANS, size=10, bold=True, color=ACCENT)
                run(p, it, font=SANS, size=10, color=TEXT, space_after=Pt(2))

        cur_y += h

        if i < len(rows) - 1:
            add_thin_rule(s, inner_x, cur_y, inner_w, height=rule_thick)
            cur_y += rule_thick


def slide_closing(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s, BG_PANEL)
    add_brand_strip(s)
    box_w = Inches(10)
    box_h = Inches(3.5)
    x = (SLIDE_W - box_w) / 2
    y = (SLIDE_H - box_h) / 2
    tf = add_text_box(s, x, y, box_w, box_h, anchor=MSO_ANCHOR.MIDDLE)
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run(p, "Thank you.", font=SERIF, size=CLOSING_SIZE, color=HEADLINE,
        align=PP_ALIGN.CENTER)


# ───────────────────────────────────────────────────────────────────────────
# Build
# ───────────────────────────────────────────────────────────────────────────

def main():
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    slide_hero(prs)
    slide_vision(prs)
    slide_commitments(prs)
    slide_why_now(prs)
    slide_who(prs)
    slide_experience(prs)
    slide_unlock(prs)
    slide_how(prs)
    slide_phase1(prs)
    slide_closing(prs)

    out = Path(__file__).parent / "sf-civic-art-map.pptx"
    prs.save(str(out))
    print(f"wrote {out} ({out.stat().st_size:,} bytes, {len(prs.slides)} slides)")


if __name__ == "__main__":
    main()
